from django.shortcuts import render
import fitz  # PyMuPDF
from django.http import HttpResponse
from PyPDF2 import PdfMerger
from pdf2docx import Converter
import os
from django.core.files.storage import default_storage, FileSystemStorage
from django.conf import settings
import pdfplumber
import pandas as pd
import uuid
from pathlib import Path
from openpyxl import Workbook
import tempfile
from pdf2image import convert_from_bytes
from pptx import Presentation
from pptx.util import Inches

def dashboard(request):
    return render(request, "dashboard.html")

def pdf_text_extractor(request):
    text = None
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        pdf_file = request.FILES['pdf_file']
        doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
        text = "\n".join([page.get_text() for page in doc])
        doc.close()
    return render(request, "pdf_text_extractor.html", {'extracted_text': text})


def pdf_merger(request):
    if request.method == 'POST' and request.FILES.getlist('pdf_files'):
        pdf_files = request.FILES.getlist('pdf_files')
        merger = PdfMerger()

        for pdf in pdf_files:
            merger.append(pdf)

        response = HttpResponse(content_type='application/pdf')
        response['Content-Disposition'] = 'attachment; filename="merged.pdf"'
        merger.write(response)
        merger.close()
        return response

    return render(request, 'pdf_merger.html')


def pdf_to_word(request):
    converted = False
    docx_file_url = None

    if request.method == 'POST' and request.FILES.get('pdf_file'):
        pdf_file = request.FILES['pdf_file']
        pdf_path = default_storage.save('temp/' + pdf_file.name, pdf_file)
        pdf_full_path = os.path.join(settings.MEDIA_ROOT, pdf_path)

        docx_output = pdf_full_path.replace('.pdf', '.docx')
        cv = Converter(pdf_full_path)
        cv.convert(docx_output)
        cv.close()

        docx_file_url = settings.MEDIA_URL + pdf_path.replace('.pdf', '.docx')
        converted = True

    return render(request, 'pdf_to_word.html', {
        'converted': converted,
        'docx_file_url': docx_file_url
    })


def pdf_to_excel(request):
    excel_url = None

    if request.method == 'POST' and request.FILES.get('pdf_file'):
        pdf_file = request.FILES['pdf_file']
        unique_id = str(uuid.uuid4())
        pdf_path = Path(settings.MEDIA_ROOT) / f"{unique_id}.pdf"
        
        with open(pdf_path, 'wb+') as f:
            for chunk in pdf_file.chunks():
                f.write(chunk)

        # Start converting PDF to Excel
        wb = Workbook()
        ws = wb.active

        with pdfplumber.open(pdf_path) as pdf:
            row_num = 1
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    for line in text.split('\n'):
                        ws.cell(row=row_num, column=1).value = line
                        row_num += 1

        # Save the Excel file
        excel_filename = f"{unique_id}.xlsx"
        excel_path = Path(settings.MEDIA_ROOT) / excel_filename
        wb.save(excel_path)
        excel_url = settings.MEDIA_URL + excel_filename

    return render(request, 'pdf_to_excel.html', {'excel_url': excel_url})


def pdf_to_ppt(request):
    if request.method == 'POST' and request.FILES.get('pdf_file'):
        pdf_file = request.FILES['pdf_file']
        
        # Convert PDF to images
        images = convert_from_bytes(pdf_file.read(), poppler_path="C:/poppler-xx/bin")

        # Create presentation
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]  # blank layout

        for img in images:
            slide = prs.slides.add_slide(blank_slide_layout)

            # Save temp image
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                img_path = tmp.name
                img.save(img_path, 'PNG')

            # Slide size
            prs_width = prs.slide_width
            prs_height = prs.slide_height

            # Add image to slide
            slide.shapes.add_picture(img_path, 0, 0, width=prs_width, height=prs_height)

            # Clean up temp image
            os.remove(img_path)

        # Save PPT to memory
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_ppt:
            prs.save(tmp_ppt.name)
            tmp_ppt.seek(0)
            pptx_data = tmp_ppt.read()

        response = HttpResponse(pptx_data, content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = 'attachment; filename="converted.pptx"'
        return response

    return render(request, 'pdf_to_ppt.html')
